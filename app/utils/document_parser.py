import PyPDF2
import docx
import pptx
from typing import List

def parse_pdf(file_path: str) -> str:
    """Parse PDF file and return text content"""
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def parse_docx(file_path: str) -> str:
    """Parse DOCX file and return text content"""
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def parse_txt(file_path: str) -> str:
    """Parse TXT file and return text content"""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def parse_pptx(file_path: str) -> str:
    """Parse PPTX file and return text content"""
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_document(file_path: str, file_type: str) -> str:
    """Parse document based on file type"""
    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "txt": parse_txt,
        "pptx": parse_pptx
    }
    
    extension = file_type.lower()
    if extension in parsers:
        return parsers[extension](file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")